import logging
from pathlib import Path
import shutil
from fsf_core.exceptions import HandlerError

log = logging.getLogger('fsf')

SUPPORTED_EXTENSIONS = {".docx", ".xlsx", ".pptx"}

class OfficeHandler:
    @staticmethod
    def can_handle(filepath: Path) -> bool:
        return filepath.suffix.lower() in SUPPORTED_EXTENSIONS
    
    @staticmethod
    def view_metadata(filepath: Path) -> dict:
        result = {'core': {}}
        ext = filepath.suffix.lower()
        try:
            if ext == '.docx':
                import docx
                doc = docx.Document(filepath)
                result['core'] = OfficeHandler._extract_props(doc.core_properties)
            elif ext == '.xlsx':
                import openpyxl
                wb = openpyxl.load_workbook(filepath, read_only=True)
                result['core'] = OfficeHandler._extract_props(wb.properties)
            elif ext == '.pptx':
                try:
                    import pptx
                    prs = pptx.Presentation(filepath)
                    result['core'] = OfficeHandler._extract_props(prs.core_properties)
                except ImportError:
                    pass
        except Exception as e:
            result['core']['Error'] = str(e)
        return result
    
    @staticmethod
    def clean_metadata(filepath: Path, output: Path | None = None) -> Path:
        out_path = output or filepath
        if out_path != filepath:
            out_path.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(filepath, out_path)
            
        ext = out_path.suffix.lower()
        try:
            if ext == '.docx':
                import docx
                doc = docx.Document(out_path)
                OfficeHandler._clear_props(doc.core_properties)
                doc.save(out_path)
            elif ext == '.xlsx':
                import openpyxl
                wb = openpyxl.load_workbook(out_path)
                OfficeHandler._clear_props(wb.properties)
                wb.save(out_path)
            elif ext == '.pptx':
                try:
                    import pptx
                    prs = pptx.Presentation(out_path)
                    OfficeHandler._clear_props(prs.core_properties)
                    prs.save(out_path)
                except ImportError:
                    pass
        except Exception:
            pass
        return out_path
    
    @staticmethod
    def spoof_metadata(filepath: Path, data: dict, output: Path | None = None) -> Path:
        out_path = output or filepath
        if out_path != filepath:
            out_path.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(filepath, out_path)
            
        ext = out_path.suffix.lower()
        try:
            if ext == '.docx':
                import docx
                doc = docx.Document(out_path)
                OfficeHandler._set_props(doc.core_properties, data)
                doc.save(out_path)
            elif ext == '.xlsx':
                import openpyxl
                wb = openpyxl.load_workbook(out_path)
                OfficeHandler._set_props(wb.properties, data)
                wb.save(out_path)
            elif ext == '.pptx':
                try:
                    import pptx
                    prs = pptx.Presentation(out_path)
                    OfficeHandler._set_props(prs.core_properties, data)
                    prs.save(out_path)
                except ImportError:
                    pass
        except Exception:
            pass
        return out_path
    
    @staticmethod
    def clone_metadata(source: Path, target: Path, output: Path | None = None) -> Path:
        out_path = output or target
        if out_path != target:
            out_path.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(target, out_path)
            
        src_ext = source.suffix.lower()
        tgt_ext = target.suffix.lower()
        
        if src_ext != tgt_ext:
            return out_path
            
        try:
            if tgt_ext == '.docx':
                import docx
                src_doc = docx.Document(source)
                tgt_doc = docx.Document(out_path)
                OfficeHandler._copy_props(src_doc.core_properties, tgt_doc.core_properties)
                tgt_doc.save(out_path)
            elif tgt_ext == '.xlsx':
                import openpyxl
                src_wb = openpyxl.load_workbook(source, read_only=True)
                tgt_wb = openpyxl.load_workbook(out_path)
                OfficeHandler._copy_props(src_wb.properties, tgt_wb.properties)
                tgt_wb.save(out_path)
            elif tgt_ext == '.pptx':
                try:
                    import pptx
                    src_prs = pptx.Presentation(source)
                    tgt_prs = pptx.Presentation(out_path)
                    OfficeHandler._copy_props(src_prs.core_properties, tgt_prs.core_properties)
                    tgt_prs.save(out_path)
                except ImportError:
                    pass
        except Exception:
            pass
        return out_path

    @staticmethod
    def _extract_props(props) -> dict:
        result = {}
        fields = ['author', 'creator', 'title', 'subject', 'description', 'keywords', 'category', 'created', 'modified', 'last_modified_by', 'lastModifiedBy', 'revision', 'language']
        for field in fields:
            try:
                val = getattr(props, field, None)
                if val:
                    if field == 'creator':
                        result['author'] = str(val)
                    elif field == 'lastModifiedBy':
                        result['last_modified_by'] = str(val)
                    else:
                        result[field] = str(val)
            except Exception:
                pass
        return result
        
    @staticmethod
    def _clear_props(props):
        fields = ['author', 'creator', 'title', 'subject', 'description', 'keywords', 'category', 'last_modified_by', 'lastModifiedBy', 'revision', 'language']
        for field in fields:
            try:
                if hasattr(props, field):
                    setattr(props, field, "")
            except Exception:
                pass
                
        for field in ['created', 'modified']:
            try:
                if hasattr(props, field):
                    setattr(props, field, None)
            except Exception:
                pass

    @staticmethod
    def _set_props(props, data: dict):
        mapping = {
            'author': ['author', 'creator'],
            'last_modified_by': ['last_modified_by', 'lastModifiedBy'],
        }
        for k, v in data.items():
            fields_to_set = mapping.get(k, [k])
            for field in fields_to_set:
                if hasattr(props, field):
                    try:
                        setattr(props, field, v)
                    except Exception:
                        pass

    @staticmethod
    def _copy_props(src_props, tgt_props):
        fields = ['author', 'creator', 'title', 'subject', 'description', 'keywords', 'category', 'created', 'modified', 'last_modified_by', 'lastModifiedBy', 'revision', 'language']
        for field in fields:
            try:
                val = getattr(src_props, field, None)
                if hasattr(tgt_props, field):
                    setattr(tgt_props, field, val)
            except Exception:
                pass
